from dotenv import load_dotenv
import os
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
# from langchain.embeddings.openai import OpenAIEmbeddings
# from langchain.vectorstores import FAISS
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from pptx import Presentation
from langchain.callbacks import get_openai_callback
load_dotenv()

os.getenv("OPENAI_API_KEY")

def main():
    st.set_page_config(page_title="Ask your pdf")
    st.header("Ask your pdf")

    pdf = st.file_uploader("Upload your pdf", type=["pdf", "pptx"])
    if pdf is not None:
        # st.write(pdf)
        if pdf.type == "application/pdf":
            pdf_reader = PdfReader(pdf)
            text=""
            for page in pdf_reader.pages:
                text += page.extract_text()

        elif pdf.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            presentation = Presentation(pdf)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        # st.write(text)

        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(text)
        # st.write(len(chunks))
        # st.write(chunks)


        embeddings = OpenAIEmbeddings()
        knowledge_base = FAISS.from_texts(chunks, embeddings)
        # st.write(knowledge_base)

        user_question = st.text_input("Ask a question about your pdf")
        if user_question:
            doc = knowledge_base.similarity_search(user_question)
            # st.write(doc)

            chain = load_qa_chain(OpenAI(), chain_type="stuff")
            with get_openai_callback() as cb:
                response = chain.run(input_documents=doc, question=user_question)
                print(cb)
            st.write(response)



if __name__ == "__main__":
    main()